import PyPDF2
import docx
from pptx import Presentation


def read_pdf(file_path):
    print(f"Reading PDF: {file_path}")
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        return ' '.join([page.extract_text() for page in reader.pages])


def read_docx(file_path):
    print(f"Reading DOCX: {file_path}")
    doc = docx.Document(file_path)
    return ' '.join([para.text for para in doc.paragraphs])


def read_pptx(file_path):
    print(f"Reading PPTX: {file_path}")
    prs = Presentation(file_path)
    return ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])


def chunk_text(text, chunk_size=500, overlap=50):
    print(f"Chunking text of length {len(text)} with chunk size {chunk_size} and overlap {overlap}")
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = ' '.join(words[i:i + chunk_size])
        chunks.append(chunk)
    print(f"Created {len(chunks)} chunks")
    return chunks


def read_document_chunk(file_path, chunk_id):
    print(f"Reading document chunk: {file_path}, chunk_id: {chunk_id}")
    content = ""
    if file_path.endswith('.pdf'):
        content = read_pdf(file_path)
    elif file_path.endswith('.docx'):
        content = read_docx(file_path)
    elif file_path.endswith('.txt'):
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

    chunks = chunk_text(content)
    return chunks[chunk_id] if chunk_id < len(chunks) else ""
